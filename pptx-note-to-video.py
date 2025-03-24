import yaml
import re
import os
import azure.cognitiveservices.speech as speechsdk
import logging
import uuid
import requests
import json
import datetime
import time
import argparse
from pptx import Presentation
from openai import AzureOpenAI
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Cm, Pt
from lxml import etree

speechsdk._log_level = speechsdk.LogLevel.Error

load_dotenv()

client = AzureOpenAI(azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"), api_version="2024-02-15-preview", api_key=os.getenv("AZURE_OPENAI_API_KEY"))
speech_config = speechsdk.SpeechConfig(subscription=os.getenv("SPEECH_KEY"), region=os.getenv("SPEECH_REGION"))
speech_config.set_speech_synthesis_output_format(speechsdk.SpeechSynthesisOutputFormat.Riff24Khz16BitMonoPcm)  

def generate_audio(transcript: str, audio_filename: str):
    file_config = speechsdk.audio.AudioOutputConfig(filename=audio_filename)
    speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=file_config)  

    result = speech_synthesizer.speak_text_async(transcript).get()
    return result

def generate_video(transcript, mp4_filename):
    job_id = str(uuid.uuid4())
    download_url = None

    url = f'https://{os.getenv("SPEECH_REGION")}.api.cognitive.microsoft.com/avatar/batchsyntheses/{job_id}?api-version=2024-04-15-preview'

    header = {
        'Content-Type': 'application/json',
        'Ocp-Apim-Subscription-Key': os.getenv("SPEECH_KEY")
    }

    payload = {
        'synthesisConfig': {
            "voice": 'en-US-AvaMultilingualNeural',
        },
        'customVoices': {
            # "YOUR_CUSTOM_VOICE_NAME": "YOUR_CUSTOM_VOICE_ID"
        },
        "inputKind": "plainText",
        "inputs": [
            {
                "content": transcript,
            },
        ],
        "avatarConfig":
            {
                "customized": False, # set to True if you want to use customized avatar
                "talkingAvatarCharacter": 'Lisa',  # talking avatar character
                "talkingAvatarStyle": 'technical-sitting',  # talking avatar style, required for prebuilt avatar, optional for custom avatar
                "videoFormat": "mp4",
                "videoCodec": "h264",
                "subtitleType": "external_file",
                "backgroundColor": "#FFFFFFFF", # background color in RGBA format, default is white; can be set to 'transparent' for transparent background
                "videoCrop": {  "topLeft": { "x": 560, "y": 0}, "bottomRight": { "x": 1360, "y": 1079}  }
            }  
    }
    #"videoCrop": {  "topLeft": { "x": 460, "y": 0}, "bottomRight": { "x": 1460, "y": 1079}  }

    response = requests.put(url, json.dumps(payload, default=str), headers=header)
    if response.status_code < 400:
        print(f'Job ID: {response.json()["id"]}')
    else:
        print(f'- Failed to submit batch avatar job: [{response.status_code}], {response.text}')

    while True:
        status = get_synthesis(url)
        if status == 'Succeeded':
            print('- batch avatar job succeeded')
            download_url, subtitle_url = getdownloadurl(url)
            
            response = requests.get(download_url)
            with open(mp4_filename, 'wb') as file:
                print("Saving video as ", mp4_filename)
                file.write(response.content) 
            
            local_srt_url = mp4_filename.replace('.mp4', '.srt')

            response = requests.get(subtitle_url)
            with open(local_srt_url, 'wb') as file:
                file.write(response.content)   
                
            break
        elif status == 'Failed':
            print('- batch avatar job failed')
            break
        else:
            print(f'- batch avatar job is [{status}]')
            time.sleep(5)

def get_synthesis(url):
    header = {
        'Ocp-Apim-Subscription-Key': os.getenv("SPEECH_KEY")
    }

    response = requests.get(url, headers=header)
    if response.status_code < 400:
        if response.json()['status'] == 'Succeeded':
            print(f'Download URL: {response.json()["outputs"]["result"]}')
        return response.json()['status']
    else:
        print(f'- Failed to get batch job: {response.text}')       

def getdownloadurl(url):
    header = {
        'Ocp-Apim-Subscription-Key': os.getenv("SPEECH_KEY")
    }

    response = requests.get(url, headers=header)
    if response.status_code < 400:
        print('- Get batch avatar job successfully')
        if response.json()['status'] == 'Succeeded':
            return response.json()["outputs"]["result"], response.json()["outputs"]["subtitle"]
    else:
        print(f'- Failed to get batch avatar job: {response.text}')
        
def main(input_pptx, output_pptx, slide):
    presentation = Presentation(input_pptx)
    
    if slide is not None:
        slides = [presentation.slides[slide - 1]]
    else:
        slides = presentation.slides

    previous_iteration = datetime.datetime.now() - datetime.timedelta(seconds=60)
    
    for slide in slides:
        print(f"Slide {slide.slide_id}")
        
        notes_part = slide.notes_slide
        transcript = notes_part.notes_text_frame.text
        
        if not transcript.strip():
            print("- Transcript is empty, skipping slide...")
            continue

        # print("- Generating audio...")
        # audio_filename = f"./audio/audio-{slide.slide_id}.wav"
        # generate_audio(transcript, audio_filename)
        # slide.shapes.add_movie(audio_filename, 0, 0, 1, 1, mime_type="audio/mpeg")
        
        now = datetime.datetime.now()
        seconds_since_last_iteration = (now - previous_iteration).total_seconds()
        if seconds_since_last_iteration < 35:
            print(f"- Waiting {35 - seconds_since_last_iteration} seconds before submitting video batch creation request...")
            time.sleep(35 - seconds_since_last_iteration)
        
        previous_iteration = datetime.datetime.now()
        
        print("- Generating video...")
        mp4_filename = f"./video/video-{slide.slide_id}.mp4"
        generate_video(transcript, mp4_filename)
        width = Inches(5.56)
        height = Inches(7.5)
        top = Inches(0)
        left = Inches(7.77)
        
        movie = slide.shapes.add_movie(mp4_filename, left, top, width, height, poster_frame_image=None, mime_type='video/mp4')
        # Send the movie to the back
        slide.shapes._spTree.remove(movie._element)
        slide.shapes._spTree.insert(2, movie._element)

        print("Saving presentation...")
        presentation.save(output_pptx)
        
        
    
if __name__ == "__main__":    
    parser = argparse.ArgumentParser(prog="pptx-note-to-video", description="Generate avatar video from PowerPoint notes.")
    parser.add_argument("--input_pptx", type=str, help="Path to the powerpoint file", required=True)
    parser.add_argument("--slide", type=int, help="Slide number (omit if you want to generate avatar video on each slide)", required=False)
    parser.add_argument("--output_pptx", type=str, help="Path to the powerpoint file", required=True)
    parser.print_help()
    args = parser.parse_args()
    
    print("")
    
    main(args.input_pptx, args.output_pptx, args.slide)