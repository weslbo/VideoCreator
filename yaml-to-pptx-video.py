import yaml
import re
import os
import azure.cognitiveservices.speech as speechsdk
import logging
import uuid
import requests
import json
import datetime
import time
import argparse
from pptx import Presentation
from openai import AzureOpenAI
from dotenv import load_dotenv
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Cm, Pt
from lxml import etree

speechsdk._log_level = speechsdk.LogLevel.Error

load_dotenv()

client = AzureOpenAI(azure_endpoint=os.getenv("AZURE_OPENAI_ENDPOINT"), api_version="2024-02-15-preview", api_key=os.getenv("AZURE_OPENAI_API_KEY"))
speech_config = speechsdk.SpeechConfig(subscription=os.getenv("SPEECH_KEY"), region=os.getenv("SPEECH_REGION"))
speech_config.set_speech_synthesis_output_format(speechsdk.SpeechSynthesisOutputFormat.Riff24Khz16BitMonoPcm)  

def read_yml_file(file_path):
    yml_dir = os.path.dirname(file_path)
    with open(file_path, 'r') as file:
        data = yaml.safe_load(file)
        title = data.get('title')
        uid = data.get('uid')
        
        content_path = re.search(r'\((.*?)\)', data.get('content')).group(1)   #.strip()[10:-1]  # Extracting the path from the include directive
        content_path = os.path.join(yml_dir, content_path)
        return title, uid, content_path

def read_sections_file(content_path, title):
    with open(content_path, 'r') as file:
        markdown_text = file.read()
        markdown_text = "# " + title + "\n" + markdown_text
        
        # Regex to match headers
        header_regex = re.compile(r"^(#{1,6})\s+(.*)", re.MULTILINE)

        # Find all headers with positions
        headers = [(match.start(), len(match.group(1)), match.group(2)) for match in header_regex.finditer(markdown_text)]

        # Extract sections based on headers
        sections = []
        for i in range(len(headers)):
            start_idx = headers[i][0]
            end_idx = headers[i+1][0] if i + 1 < len(headers) else len(markdown_text)
            
            header_text = headers[i][2]   # Header text
            content = markdown_text[start_idx:end_idx].split("\n", 1)[1].strip() if "\n" in markdown_text[start_idx:end_idx] else ""
            
            sections.append({"title": header_text, "content": content})

        return sections
    
def generate_bullets(content):
    prompt = f"Generate a list of bullet points based on the following content:\n\n{content}"
    
    message_text = [
        {"role":"system","content":"""
            Create a list of bullet points based on the content provided. 
            - Minimal 3 items in the list.
            - Maximum 8 items in the ist
            - Make sure to be concise, complete and clear. 
            - Each bullet points should be maximum 5 words long.
         """},
        {"role":"user","content":prompt}
    ]
        
    completion = client.chat.completions.create(
        model="gpt-4o",
        messages = message_text,
        temperature=0.1,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None
    )    
    
    output = completion.choices[0].message.content
    return output

def generate_speakertranscript(content, mode):
    prompt = f"Generate the speaker transctip:\n\n{content}"
    
    if mode == "intro":
        additional_prompt = "Welcome the audience. At the end, you do not have to say goodbye or thank the audience."
    elif mode == "conclusion":
        additional_prompt = "Do not include an introducion or greeting, but at the end, thank the audience for watching the video."
    else:
        additional_prompt = "Do not include any greetings or introductions."
    
    message_text = [
        {"role":"system","content":f"""
            Create a speaker transcript based on the content provided. 
            - You are a professional instructor and you are giving a presentation.
            - Do not include any bullet points or lists. 
            - Make sure to be concise and clear.
            - Use natural language and avoid markdown, code fragments, or any other formatting.
            - File extensions should be in uppercase (like .MP4, .PDF, etc.)
            - This will be used for a video, so make sure to use a friendly and engaging tone.
            - {additional_prompt}
            """},
        {"role":"user","content":prompt}
    ]
        
    completion = client.chat.completions.create(
        model="gpt-4o",
        messages = message_text,
        temperature=0.1,
        top_p=0.95,
        frequency_penalty=0,
        presence_penalty=0,
        stop=None
    )    
    
    output = completion.choices[0].message.content
    return output

def generate_audio(transcript: str, audio_filename: str):
    file_config = speechsdk.audio.AudioOutputConfig(filename=audio_filename)
    speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=file_config)  

    result = speech_synthesizer.speak_text_async(transcript).get()
    return result

def generate_video(transcript, mp4_filename):
    job_id = str(uuid.uuid4())
    download_url = None

    url = f'https://{os.getenv("SPEECH_REGION")}.api.cognitive.microsoft.com/avatar/batchsyntheses/{job_id}?api-version=2024-04-15-preview'

    header = {
        'Content-Type': 'application/json',
        'Ocp-Apim-Subscription-Key': os.getenv("SPEECH_KEY")
    }

    payload = {
        'synthesisConfig': {
            "voice": 'en-US-AvaMultilingualNeural',
        },
        'customVoices': {
            # "YOUR_CUSTOM_VOICE_NAME": "YOUR_CUSTOM_VOICE_ID"
        },
        "inputKind": "plainText",
        "inputs": [
            {
                "content": transcript,
            },
        ],
        "avatarConfig":
            {
                "customized": False, # set to True if you want to use customized avatar
                "talkingAvatarCharacter": 'Lisa',  # talking avatar character
                "talkingAvatarStyle": 'technical-sitting',  # talking avatar style, required for prebuilt avatar, optional for custom avatar
                "videoFormat": "mp4",
                "videoCodec": "h264",
                "subtitleType": "external_file",
                "backgroundColor": "#FFFFFFFF", # background color in RGBA format, default is white; can be set to 'transparent' for transparent background
                "videoCrop": {  "topLeft": { "x": 560, "y": 0}, "bottomRight": { "x": 1360, "y": 1079}  }
            }  
    }
    #"videoCrop": {  "topLeft": { "x": 460, "y": 0}, "bottomRight": { "x": 1460, "y": 1079}  }

    response = requests.put(url, json.dumps(payload, default=str), headers=header)
    if response.status_code < 400:
        print(f'Job ID: {response.json()["id"]}')
    else:
        print(f'- Failed to submit batch avatar job: [{response.status_code}], {response.text}')

    while True:
        status = get_synthesis(url)
        if status == 'Succeeded':
            print('- batch avatar job succeeded')
            download_url, subtitle_url = getdownloadurl(url)
            
            response = requests.get(download_url)
            with open(mp4_filename, 'wb') as file:
                file.write(response.content) 
            
            local_srt_url = mp4_filename.replace('.mp4', '.srt')

            response = requests.get(subtitle_url)
            with open(local_srt_url, 'wb') as file:
                file.write(response.content)   
                
            break
        elif status == 'Failed':
            print('- batch avatar job failed')
            break
        else:
            print(f'- batch avatar job is [{status}]')
            time.sleep(5)

def get_synthesis(url):
    header = {
        'Ocp-Apim-Subscription-Key': os.getenv("SPEECH_KEY")
    }

    response = requests.get(url, headers=header)
    if response.status_code < 400:
        if response.json()['status'] == 'Succeeded':
            print(f'Download URL: {response.json()["outputs"]["result"]}')
        return response.json()['status']
    else:
        print(f'- Failed to get batch job: {response.text}')       

def getdownloadurl(url):
    header = {
        'Ocp-Apim-Subscription-Key': os.getenv("SPEECH_KEY")
    }

    response = requests.get(url, headers=header)
    if response.status_code < 400:
        print('- Get batch avatar job successfully')
        if response.json()['status'] == 'Succeeded':
            return response.json()["outputs"]["result"], response.json()["outputs"]["subtitle"]
    else:
        print(f'- Failed to get batch avatar job: {response.text}')

def main(yml_file):
    title, uid, content_path = read_yml_file(yml_file)
    sections = read_sections_file(content_path, title)
    presentation = Presentation("template.pptx")
    
    for section in sections:
        mode = "intro" if "intro" in yml_file else "content"
        
        slide_layout = presentation.slide_layouts[3]
        print(f"{mode.upper()}: {section["title"]}")
            
        slide = presentation.slides.add_slide(slide_layout)
        title_element = slide.shapes.title
        title_element.text = section["title"].strip() 
        
        print("- Generating bulleted list")
        content_placeholder = next(shape for shape in slide.placeholders if shape.name == "Content Placeholder 2")
        text_box = content_placeholder.text_frame
        bullets = generate_bullets(section["content"])
        for bullet in bullets.split('\n'):
            p = text_box.add_paragraph()
            p.text = bullet.replace("- ", "").strip()
    
        print("- Generating speaker transcript")
        speaker_transcript = generate_speakertranscript(section["content"], mode)
        notes_part = slide.notes_slide
        notes_part.notes_text_frame.text = speaker_transcript 
        
        print("- Generating audio")
        audio_filename = f"./audio/audio-{sections.index(section)}.wav"
        generate_audio(speaker_transcript, audio_filename)
        audio = slide.shapes.add_movie(audio_filename, 0, 0, 1, 1, mime_type="audio/mpeg")
        
        print("- Generating video")
        mp4_filename = f"./video/video-{sections.index(section)}.mp4"
        generate_video(speaker_transcript, mp4_filename)
        
        width = Inches(5.56)
        height = Inches(7.5)
        top = Inches(0)
        left = Inches(7.77)
        
        movie = slide.shapes.add_movie(mp4_filename, left, top, width, height, poster_frame_image=None, mime_type='video/mp4')
        
        # Send the movie to the back
        slide.shapes._spTree.remove(movie._element)
        slide.shapes._spTree.insert(2, movie._element)
                
    presentation.save(f"output/{uid}.pptx")
    
if __name__ == "__main__":
    parser = argparse.ArgumentParser(prog="yaml-to-pptx-video", description="Generate powerpoint video from a module YAML file.")
    parser.add_argument("--yml_file", type=str, help="Path to the YAML file.", required=True)
    parser.print_help()
    args = parser.parse_args()
    
    print("")
    
    main(args.yml_file)